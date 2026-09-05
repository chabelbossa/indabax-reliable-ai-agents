"""Build editable teaching slides using the official IndabaX theme and branding."""
from copy import deepcopy
from pathlib import Path
import os
import subprocess

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from deck_content import FR, EN

ROOT = Path(__file__).resolve().parents[1]
OFFICIAL_TEMPLATE = ROOT.parent / "IndabaX Bénin — Speaker Deck Template.pptx"
REVIEWED_TEMPLATE = Path("/Users/user/Downloads/indabax-reliable-ai-agents.pptx")
SOURCE = Path(os.environ.get("INDABAX_DECK_TEMPLATE", str(
    REVIEWED_TEMPLATE if REVIEWED_TEMPLATE.exists() else OFFICIAL_TEMPLATE)))
GREEN, INK, MUTED = "008751", "111713", "53645A"
PALE, YELLOW, RED = "EFF6F1", "FDC116", "BE2440"


def box(slide, x, y, w, h, fill=PALE, rounded=True):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if rounded:
        shape.adjustments[0] = 0.08
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill)
    shape.line.fill.background()
    return shape


def text(slide, value, x, y, w, h, size=26, color=INK, bold=False, mono=False, align=None, url=None):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    for i, line_text in enumerate(value.split("\n")):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.space_after = Pt(7 if not mono else 3)
        p.line_spacing = 1.13
        if align is not None:
            p.alignment = align
        r = p.add_run()
        r.text = line_text
        r.font.name = "DejaVu Sans Mono" if mono else "Poppins"
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = RGBColor.from_string(color)
        if url:
            r.hyperlink.address = url
    return shape


def card(slide, x, y, w, h, title, body, accent=GREEN, dark=False, size=25):
    box(slide, x, y, w, h, GREEN if dark else PALE)
    text(slide, title, x+.35, y+.3, w-.7, 1.0, 30, "FFFFFF" if dark else accent, True)
    text(slide, body, x+.35, y+1.45, w-.7, h-1.65, size, "FFFFFF" if dark else INK)


def arrow(slide, x, y, w=.6, h=.4, direction="right"):
    kind = {"right": MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
            "down": MSO_AUTO_SHAPE_TYPE.DOWN_ARROW,
            "up": MSO_AUTO_SHAPE_TYPE.UP_ARROW}[direction]
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(GREEN)
    shape.line.fill.background()


def line(slide, x1, y1, x2, y2):
    s = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    s.line.color.rgb = RGBColor.from_string(GREEN)
    s.line.width = Pt(3)


def render(slide, item):
    kind = item["kind"]
    text(slide, item["label"], 1.25, .78, 17.5, .4, 19, GREEN, True)
    text(slide, item["title"], 1.25, 1.45, 17.5, 1.5, 43, bold=True)
    if kind == "agenda":
        for i, (title, detail, duration) in enumerate(item["rows"]):
            y = 3.15+i*1.42
            text(slide, f"0{i+1}", 1.25, y, 1, .65, 31, GREEN, True)
            text(slide, title, 2.65, y, 12.5, .65, 29, bold=True)
            text(slide, detail, 2.65, y+.6, 12.5, .6, 20, MUTED)
            text(slide, duration, 16.6, y+.15, 2.15, .6, 24, GREEN, True, align=PP_ALIGN.RIGHT)
    elif kind == "mission":
        card(slide, 1.25, 3.15, 6.0, 5.65, item["alert_title"], item["alert"], dark=True, size=26)
        text(slide, item["story"], 8.0, 3.15, 10.7, 2.4, 29)
        text(slide, item["role"], 8.0, 5.9, 10.7, .6, 27, GREEN, True)
        text(slide, item["deliverable"], 8.0, 6.75, 10.7, 1.9, 26)
    elif kind == "cards":
        count = len(item["cards"])
        width = (17.5-.45*(count-1))/count
        for i, (title, body) in enumerate(item["cards"]):
            card(slide, 1.25+i*(width+.45), 3.15, width, 4.95, title, body,
                 accent=[GREEN, INK, RED][i % 3], size=item.get("body_size",25))
        if item.get("takeaway"):
            text(slide, item["takeaway"], 1.25, 8.45, 17.5, 1.0, 25, GREEN, True)
    elif kind == "tool_call":
        for x, label, code, caption in zip((1.25,10.25), item["labels"], item["code"], item["captions"]):
            text(slide, label, x, 3.05, 8.5, .65, 25, GREEN, True)
            box(slide, x, 3.95, 8.5, 3.6, INK)
            text(slide, code, x+.3, 4.25, 7.9, 3.0, 22, "FFFFFF", mono=True)
            text(slide, caption, x, 7.95, 8.5, 1.35, 23)
    elif kind == "loop":
        for x, (title, body) in zip((1.25,7.3,13.35), item["nodes"]):
            card(slide, x, 3.5, 5.4, 2.75, title, body, size=23)
        arrow(slide,6.72,4.6,.48,.4)
        arrow(slide,12.77,4.6,.48,.4)
        arrow(slide,15.75,6.5,.45,.7,"down")
        box(slide,7.3,7.35,11.45,1.05,GREEN)
        text(slide,item["observation"],7.65,7.56,10.75,.7,25,"FFFFFF",True)
        line(slide,7.3,7.88,3.95,7.88)
        line(slide,3.95,7.88,3.95,6.95)
        arrow(slide,3.72,6.45,.46,.55,"up")
        text(slide,item["return_label"],1.25,8.35,5.4,.8,22,GREEN)
        text(slide,item["stop"],1.25,2.85,17.5,.55,23,MUTED)
        text(slide,item["exit"],1.25,9.25,17.5,.7,23,GREEN,True)
    elif kind == "table":
        widths = item.get("widths",[5.8,6.9,4.8])
        y = 3.15
        box(slide,1.25,y,17.5,.8,GREEN,rounded=False)
        x=1.25
        for label,w in zip(item["headers"],widths):
            text(slide,label,x+.23,y+.16,w-.46,.5,22,"FFFFFF",True)
            x+=w
        for i,row in enumerate(item["rows"]):
            y=3.95+i*item.get("row_height",1.0)
            box(slide,1.25,y,17.5,item.get("row_height",1.0),PALE if i%2==0 else "FAFBFA",rounded=False)
            x=1.25
            for j,(value,w) in enumerate(zip(row,widths)):
                text(slide,value,x+.23,y+.18,w-.46,item.get("row_height",1.0)-.24,
                     item.get("table_size",23),GREEN if value=="PASS" else INK,
                     bold=(j==0 or value=="PASS"),mono=(j==item.get("mono_column",-1)))
                x+=w
        if item.get("takeaway"):
            text(slide,item["takeaway"],1.25,9.25,17.5,.6,23,GREEN,True)
    elif kind == "checkpoint":
        text(slide,item["duration"],16.8,.8,1.95,.55,22,GREEN,True,align=PP_ALIGN.RIGHT)
        for i,(label,body) in enumerate(item["steps"]):
            y=3.05+i*1.55
            box(slide,1.25,y,.7,.7,GREEN)
            text(slide,str(i+1),1.25,y+.1,.7,.5,25,"FFFFFF",True,align=PP_ALIGN.CENTER)
            text(slide,label,2.4,y,15.8,.6,27,bold=True)
            text(slide,body,2.4,y+.66,15.8,.75,24)
        box(slide,1.25,8.0,17.5,1.55,PALE)
        text(slide,item["success"],1.65,8.3,16.7,1.0,25,GREEN,True)
    elif kind == "failure":
        card(slide,1.25,3.1,8.45,4.5,item["left_title"],item["left_body"],accent=RED,size=27)
        card(slide,10.25,3.1,8.5,4.5,item["right_title"],item["right_body"],size=27)
        box(slide,1.25,8.05,17.5,1.35,GREEN)
        text(slide,item["decision"],1.65,8.32,16.7,.95,26,"FFFFFF",True)
    elif kind == "dossier":
        for x,(label,value) in zip((1.25,5.7,10.15,14.6),item["metrics"]):
            box(slide,x,3.1,4.15,1.6,PALE)
            text(slide,label,x+.25,3.33,3.65,.45,18,MUTED)
            text(slide,value,x+.25,3.92,3.65,.65,30,GREEN,True)
        box(slide,1.25,5.1,8.4,4.25,INK)
        text(slide,item["code"],1.6,5.45,7.7,3.55,22,"FFFFFF",mono=True)
        text(slide,item["body"],10.3,5.15,8.4,2.7,27)
        text(slide,item["takeaway"],10.3,8.0,8.4,1.2,25,GREEN,True)
    elif kind == "resources":
        for i,(title,description,url) in enumerate(item["links"]):
            y=3.1+i*1.9
            text(slide,title,1.25,y,17.5,.7,29,GREEN,True,url=url)
            text(slide,description,1.25,y+.7,17.5,.8,24)
    elif kind == "closing":
        text(slide,item["prompt"],1.25,3.25,11.8,2.25,35,GREEN,True)
        text(slide,item["example"],1.25,6.05,11.8,1.6,26)
        text(slide,item["speaker"],1.25,8.45,12,1.1,21,MUTED)
        slide.shapes.add_picture(str(ROOT/"slides/repository-qr.png"),Inches(14.65),Inches(3.35),Inches(3.5),Inches(3.5))
        text(slide,item["qr_label"],13.8,7.15,5.0,1.05,23,GREEN,True,align=PP_ALIGN.CENTER,
             url="https://github.com/chabelbossa/indabax-reliable-ai-agents")
    else:
        raise ValueError(f"Unknown layout: {kind}")
    if item.get("footnote"):
        text(slide,item["footnote"],1.25,9.8,17.5,.62,17,MUTED)


def build_deck(language):
    plan = FR if language == "fr" else EN
    presentation = Presentation(SOURCE)
    # Keep the cover, theme and footer, including the original crest image.
    footer_source = presentation.slides[3]
    footer = [deepcopy(footer_source.shapes[i].element) for i in (9,10)]
    footer_relations = dict(footer_source.part.rels.items())
    cover = presentation.slides[0]
    originals = list(presentation.slides)
    for index,value in plan[0]["text"].items():
        frame=cover.shapes[index].text_frame
        frame.paragraphs[0].runs[0].text=value
        for p in frame.paragraphs:
            for r in list(p.runs)[1:]:
                r.text=""
        for p in list(frame.paragraphs)[1:]:
            p.clear()
    cover.shapes[5].width=Inches(8)
    cover.shapes[5].height=Inches(.8)
    cover.shapes[6].height=Inches(.8)
    cover.shapes[4].width=Inches(17.5)
    cover.shapes[4].height=Inches(1.15)
    cover.notes_slide.notes_text_frame.text=plan[0]["notes"]
    for number,item in enumerate(plan[1:],2):
        slide = originals[number-1] if number <= len(originals) else presentation.slides.add_slide(presentation.slide_layouts[6])
        for s in list(slide.shapes):
            slide.shapes._spTree.remove(s.element)
        box(slide,0,0,20,11.25,"FFFFFF",rounded=False)
        for element in footer:
            clone = deepcopy(element)
            for child in clone.iter():
                for attr in (qn("r:embed"), qn("r:link")):
                    if child.get(attr):
                        rel = footer_relations[child.get(attr)]
                        target = rel.target_ref if rel.is_external else rel.target_part
                        child.set(attr, slide.part.relate_to(target, rel.reltype, rel.is_external))
            slide.shapes._spTree.insert_element_before(clone,"p:extLst")
        for s in list(slide.shapes)[1:3]:
            s.top+=Inches(.92)
        text(slide,f"{number:02d} / {len(plan):02d}",17.35,10.55,1.4,.4,17,MUTED,align=PP_ALIGN.RIGHT)
        render(slide,item)
        slide.notes_slide.notes_text_frame.text=item["notes"]
    output=ROOT/"slides"/f"indabax-reliable-ai-agents-{language}.pptx"
    presentation.save(output)
    return output


def convert_to_pdf(path):
    subprocess.run(["soffice","--headless","--convert-to","pdf","--outdir",str(path.parent),str(path)],
                   check=True,capture_output=True,text=True)
    output=path.with_suffix(".pdf")
    if not output.exists() or output.stat().st_mtime<path.stat().st_mtime:
        raise RuntimeError(f"PDF was not refreshed: {output}")
    return output


if __name__=="__main__":
    for language in ("fr","en"):
        deck=build_deck(language)
        print(deck)
        print(convert_to_pdf(deck))
